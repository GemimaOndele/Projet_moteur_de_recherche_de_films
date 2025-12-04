"""
Script pour mettre à jour la présentation PowerPoint avec :
- Graphiques générés
- Images prompts/réponses
- Tableaux de résultats
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

# Configuration
PROJECT_ROOT = Path(__file__).parent.parent
PRESENTATION_PATH = PROJECT_ROOT / "MovieMood_Presentation.pptx"
IMAGES_DIR = PROJECT_ROOT / "presentation_projet" / "images"

# Couleurs du thème
PRIMARY_COLOR = RGBColor(26, 26, 46)  # #1a1a2e
ACCENT_COLOR = RGBColor(233, 69, 96)  # #e94560
TEXT_COLOR = RGBColor(255, 255, 255)  # #ffffff
SECONDARY_COLOR = RGBColor(15, 52, 96)  # #0f3460

def add_slide_with_image(prs, title_text, image_path, layout_idx=6):
    """Ajoute un slide avec une image."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Titre
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    # Image
    if image_path.exists():
        left_img = Inches(0.5)
        top_img = Inches(1.2)
        width_img = Inches(9)
        height_img = Inches(6)
        slide.shapes.add_picture(str(image_path), left_img, top_img, width_img, height_img)
    else:
        # Texte d'erreur si l'image n'existe pas
        left = Inches(1)
        top = Inches(3)
        width = Inches(8)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"⚠️ Image non trouvée : {image_path.name}"
        p = tf.paragraphs[0]
        p.font.size = Pt(18)
        p.font.color.rgb = TEXT_COLOR

def add_slide_with_two_images(prs, title_text, image1_path, image2_path, layout_idx=6):
    """Ajoute un slide avec deux images côte à côte."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    
    # Fond
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Titre
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = title_text
    p = tf.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.LEFT
    
    # Image 1 (gauche)
    if image1_path.exists():
        left_img1 = Inches(0.5)
        top_img1 = Inches(1.2)
        width_img = Inches(4.5)
        height_img = Inches(6)
        slide.shapes.add_picture(str(image1_path), left_img1, top_img1, width_img, height_img)
    
    # Image 2 (droite)
    if image2_path.exists():
        left_img2 = Inches(5.5)
        top_img2 = Inches(1.2)
        width_img = Inches(4.5)
        height_img = Inches(6)
        slide.shapes.add_picture(str(image2_path), left_img2, top_img2, width_img, height_img)

def update_presentation():
    """Met à jour la présentation avec tous les éléments générés."""
    print("=" * 60)
    print("📊 MISE À JOUR DE LA PRÉSENTATION POWERPOINT")
    print("=" * 60)
    
    # Charger ou créer la présentation
    if PRESENTATION_PATH.exists():
        print(f"📂 Chargement de la présentation existante : {PRESENTATION_PATH}")
        prs = Presentation(str(PRESENTATION_PATH))
    else:
        print("⚠️  Présentation non trouvée, création d'une nouvelle...")
        prs = Presentation()
    
    # Vérifier que le dossier images existe
    if not IMAGES_DIR.exists():
        print(f"❌ Erreur : Le dossier {IMAGES_DIR} n'existe pas !")
        print("💡 Exécutez d'abord : python scripts/enrich_presentation.py")
        return
    
    print(f"\n📁 Dossier images : {IMAGES_DIR}")
    
    # Ajouter les slides avec les graphiques
    print("\n1️⃣  Ajout des graphiques...")
    
    # Slide : Distributions
    dist_path = IMAGES_DIR / "01_distributions.png"
    if dist_path.exists():
        add_slide_with_image(prs, "📊 Distribution des Données", dist_path)
        print("   ✅ Slide ajouté : Distribution des Données")
    
    # Slide : Top Genres
    genres_path = IMAGES_DIR / "02_top_genres.png"
    if genres_path.exists():
        add_slide_with_image(prs, "🎭 Top 15 Genres les Plus Représentés", genres_path)
        print("   ✅ Slide ajouté : Top Genres")
    
    # Slide : Performance Recommandations
    perf_path = IMAGES_DIR / "03_performance_recommandations.png"
    if perf_path.exists():
        add_slide_with_image(prs, "🎯 Performance des Recommandations par Émotion", perf_path)
        print("   ✅ Slide ajouté : Performance Recommandations")
    
    # Slide : Matrice de Confusion
    confusion_path = IMAGES_DIR / "04_matrice_confusion.png"
    if confusion_path.exists():
        add_slide_with_image(prs, "📊 Matrice de Confusion : Genres × Émotions", confusion_path)
        print("   ✅ Slide ajouté : Matrice de Confusion")
    
    # Ajouter les slides avec les tableaux
    print("\n2️⃣  Ajout des tableaux...")
    
    # Slide : Tableau Résumé
    resume_path = IMAGES_DIR / "05_tableau_resume.png"
    if resume_path.exists():
        add_slide_with_image(prs, "📋 Résumé des Résultats", resume_path)
        print("   ✅ Slide ajouté : Tableau Résumé")
    
    # Slide : Tableau Recommandations
    rec_table_path = IMAGES_DIR / "06_tableau_recommandations.png"
    if rec_table_path.exists():
        add_slide_with_image(prs, "📊 Tableau des Recommandations par Émotion", rec_table_path)
        print("   ✅ Slide ajouté : Tableau Recommandations")
    
    # Ajouter les slides avec les prompts/réponses
    print("\n3️⃣  Ajout des prompts et réponses...")
    
    for i in range(1, 7):
        prompt_path = IMAGES_DIR / f"prompt_{i:02d}.png"
        response_path = IMAGES_DIR / f"response_{i:02d}.png"
        
        if prompt_path.exists() and response_path.exists():
            add_slide_with_two_images(
                prs,
                f"💬 Prompt {i} : Question & Réponse",
                prompt_path,
                response_path
            )
            print(f"   ✅ Slide ajouté : Prompt {i}")
        elif prompt_path.exists():
            add_slide_with_image(prs, f"💬 Prompt {i}", prompt_path)
            print(f"   ✅ Slide ajouté : Prompt {i} (sans réponse)")
    
    # Sauvegarder la présentation
    output_path = PROJECT_ROOT / "MovieMood_Presentation_Enrichie.pptx"
    prs.save(str(output_path))
    
    print("\n" + "=" * 60)
    print("✅ PRÉSENTATION MIS À JOUR AVEC SUCCÈS !")
    print(f"📁 Fichier sauvegardé : {output_path}")
    print(f"📊 Nombre total de slides : {len(prs.slides)}")
    print("=" * 60)
    
    return output_path

if __name__ == "__main__":
    try:
        output_path = update_presentation()
        print(f"\n🎉 La présentation enrichie est prête dans : {output_path}")
        print("💡 Vous pouvez maintenant l'ouvrir dans PowerPoint !")
    except ImportError:
        print("❌ Erreur : python-pptx n'est pas installé.")
        print("💡 Installez-le avec : pip install python-pptx")
    except Exception as e:
        print(f"❌ Erreur lors de la mise à jour : {e}")
        import traceback
        traceback.print_exc()

