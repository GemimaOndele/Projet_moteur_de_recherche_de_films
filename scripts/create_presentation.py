"""
Script pour générer automatiquement la présentation PowerPoint MovieMood
Crée un fichier .pptx avec tous les slides de la présentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pathlib import Path

def create_presentation():
    """Crée la présentation PowerPoint complète."""
    
    # Créer une nouvelle présentation
    prs = Presentation()
    
    # Définir les couleurs du thème
    PRIMARY_COLOR = RGBColor(26, 26, 46)  # #1a1a2e
    ACCENT_COLOR = RGBColor(233, 69, 96)  # #e94560
    TEXT_COLOR = RGBColor(255, 255, 255)  # #ffffff
    SECONDARY_COLOR = RGBColor(15, 52, 96)  # #0f3460
    
    # ============================================
    # SLIDE 1 : Page de Titre
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Layout vide
    
    # Fond sombre
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    # Titre principal
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "🎬 MovieMood"
    p = tf.paragraphs[0]
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Sous-titre
    top = Inches(3.5)
    height = Inches(0.8)
    txBox2 = slide.shapes.add_textbox(left, top, width, height)
    tf2 = txBox2.text_frame
    tf2.text = "Plateforme Web IA de Recommandation de Films par Émotion"
    p2 = tf2.paragraphs[0]
    p2.font.size = Pt(24)
    p2.font.color.rgb = TEXT_COLOR
    p2.alignment = PP_ALIGN.CENTER
    
    # Équipe
    top = Inches(5)
    height = Inches(1.5)
    txBox3 = slide.shapes.add_textbox(left, top, width, height)
    tf3 = txBox3.text_frame
    tf3.text = "Équipe :\nGémima ONDELE POURU | Fatoumata BAH | Hector KOMBOU\n\nProjet Académique - Data Engineering"
    p3 = tf3.paragraphs[0]
    p3.font.size = Pt(18)
    p3.font.color.rgb = TEXT_COLOR
    p3.alignment = PP_ALIGN.CENTER
    
    # ============================================
    # SLIDE 2 : Problématique & Vision
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Layout titre et contenu
    title = slide.shapes.title
    title.text = "Problématique & Vision"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Le Défi :"
    p = tf.paragraphs[0]
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Comment choisir un film qui correspond à notre humeur du moment ?"'
    p.font.size = Pt(18)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nNotre Solution : 🎯"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "Plateforme web intelligente qui recommande des films basés sur :"
    p.font.size = Pt(16)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Votre émotion actuelle"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• L'analyse de sentiments des films"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Un système de scoring avancé"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\nImpact : ✅"
    p.font.size = Pt(20)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "• Économise du temps de recherche"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Améliore l'expérience cinématographique"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Personnalisation par émotion"
    p.font.size = Pt(14)
    p.level = 2
    
    # ============================================
    # SLIDE 3 : Architecture Technique
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Architecture Technique"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Vision Globale :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    architecture_text = """┌─────────────────────────────────────┐
│  Front-end Web (HTML/CSS/JS)        │
│  - Interface utilisateur            │
│  - Recherche par titre/émotion      │
└──────────────┬──────────────────────┘
               │
┌──────────────▼──────────────────────┐
│  Back-end Flask (API REST)          │
│  - /search?title=...                │
│  - /recommend?emotion=...           │
└──────────────┬──────────────────────┘
               │
┌──────────────▼──────────────────────┐
│  Modules IA & Données               │
│  - data_loading.py                  │
│  - sentiment.py (TextBlob)          │
│  - recommendation.py                │
│  - emotion_detection.py (DeepFace)  │
└──────────────┬──────────────────────┘
               │
┌──────────────▼──────────────────────┐
│  Dataset TMDB (4803 films)          │
│  - Métadonnées enrichies            │
│  - Sentiments analysés              │
└─────────────────────────────────────┘"""
    
    p = tf.add_paragraph()
    p.text = architecture_text
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    
    # ============================================
    # SLIDE 4 : Plan de Travail & Répartition
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Plan de Travail & Répartition"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Phase 1 - Données & Analyse de Sentiments"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "👤 Gémima ONDELE"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Chargement dataset TMDB"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Analyse de sentiments avec TextBlob"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Enrichissement des données"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\nPhase 2 - Moteur de Recherche"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "👤 Hector KOMBOU"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Recherche par titre"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Recommandations par similarité"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Recommandations par émotion"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Système de scoring"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\nPhase 3 - Interface Web"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_COLOR
    
    p = tf.add_paragraph()
    p.text = "👤 Fatoumata BAH"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Design et UX"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Pages HTML/CSS"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Intégration API"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Expérience utilisateur"
    p.font.size = Pt(14)
    p.level = 2
    
    # ============================================
    # SLIDE 5-6 : PHASE 1 - Prompt Initial
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 1 : Architecture du Projet"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT à l'IA :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Nous voulons créer une plateforme web IA de recommandation de films.\nNous avons un dataset TMDB avec 5000 films. Comment structurer le projet ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE DE L'IA :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Proposition d'architecture en 4 couches :"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "1. Front-end web (HTML/CSS/JS)"
    p.font.size = Pt(12)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "2. Back-end Python (Flask/FastAPI)"
    p.font.size = Pt(12)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "3. Modules IA/Données"
    p.font.size = Pt(12)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "4. Dataset TMDB"
    p.font.size = Pt(12)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\n✅ Plan de travail avec répartition des tâches"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Cahier des charges structuré"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 6 : Code Phase 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 1 : Code Généré"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE RÉSULTANT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """# data_loading.py (Gémima)
def charger_dataframe(path_csv: str) -> pd.DataFrame:
    \"\"\"Charge le fichier CSV TMDB.\"\"\"
    df = pd.read_csv(path_csv, engine="python", 
                     encoding="utf-8")
    return df

def construire_liste_films(df: pd.DataFrame):
    \"\"\"Construit une liste de films avec genres.\"\"\"
    films = []
    for _, row in df.iterrows():
        film = {
            "id": row["id"],
            "title": row["title"],
            "genres": parser_genres(row["genres"]),
            "overview": row["overview"],
            "vote_average": row["vote_average"],
        }
        films.append(film)
    return films"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 7-8 : PHASE 1 - Analyse de Sentiments
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 1 : Analyse de Sentiments"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Comment analyser le sentiment des résumés de films ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Utilisation de TextBlob"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Fonction analyser_sentiment_texte()"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Retourne : score [-1,1] + label"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 8 : Code Sentiment
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 1 : Code Analyse de Sentiments"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """# lib_projet.py (Gémima)
from textblob import TextBlob

def analyser_sentiment_texte(texte: str):
    \"\"\"Analyse le sentiment avec TextBlob.\"\"\"
    blob = TextBlob(texte)
    polarite = blob.sentiment.polarity
    
    if polarite > 0.1:
        label = "positif"
    elif polarite < -0.1:
        label = "negatif"
    else:
        label = "neutre"
    
    return float(polarite), label"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 9-10 : PHASE 2 - Système de Scoring
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 2 : Système de Scoring"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Comment créer un système de scoring pour réordonner les films ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Scoring combinant :"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Sentiment du film (normalisé 0-1)"
    p.font.size = Pt(12)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Note moyenne (normalisée 0-1)"
    p.font.size = Pt(12)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Pondérations ajustables"
    p.font.size = Pt(12)
    p.level = 2
    
    # Slide 10 : Code Scoring
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 2 : Code Système de Scoring"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """# lib_projet.py (Hector)
def calculer_score_film(film: dict, emotion_user: str) -> float:
    \"\"\"Calcule un score global pour un film.\"\"\"
    sentiment_norm = normaliser_sentiment(
        film.get("sentiment_score", 0.0)
    )
    note_norm = normaliser_note(
        film.get("vote_average", 0.0)
    )
    
    w_sentiment = 0.6  # importance du sentiment
    w_note = 0.4       # importance de la note
    
    score = w_sentiment * sentiment_norm + w_note * note_norm
    return float(score)"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 11-12 : PHASE 2 - Recommandations
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 2 : Recommandations par Émotion"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Comment mapper les émotions aux genres de films ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Dictionnaire emotion_to_genres"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Fonction recommander_par_emotion()"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Filtre par genres → Score → Trie"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 12 : Code Recommandations
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 2 : Code Recommandations"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """# recommendation.py (Hector)
emotion_to_genres = {
    "triste": ["Comedy", "Family", "Drama"],
    "stressé": ["Comedy", "Adventure", "Action"],
    "heureux": ["Romance", "Music", "Comedy"],
    # ... autres émotions
}

def recommander_par_emotion(emotion, films, n=20):
    genres_cibles = emotion_to_genres.get(emotion, [])
    candidats = []
    
    for film in films:
        if genres_cibles & film["genres"]:
            candidats.append(film)
    
    # Trier par note décroissante
    candidats.sort(key=lambda f: f["vote_average"], 
                   reverse=True)
    return candidats[:n]"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 13-14 : PHASE 3 - Interface Web
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 3 : Interface Web"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Comment créer une interface web moderne ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Flask + templates Jinja2"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Page d'accueil avec recherche"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Page de résultats avec cartes"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Design moderne CSS"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 14 : Code Interface
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "PHASE 3 : Code Interface Web"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """# app.py (Hector + Fatoumata)
@app.route("/search")
def search():
    titre = request.args.get("titre", "")
    emotion = request.args.get("emotion", "")
    
    resultats = []
    if titre:
        film = rechercher_par_titre(titre, films)
        if film:
            resultats.append(film)
    
    if emotion:
        resultats.extend(
            recommander_par_emotion(emotion, films, n=20)
        )
    
    return render_template("results.html", 
                          films=resultats)"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 15-16 : Vidéo de Fond
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Amélioration : Vidéo de Fond"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"La vidéo ne se lance pas sur Chrome/Edge. Comment corriger ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Problème : Autoplay bloqué"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Solution : YouTube Iframe API"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Bouton pour activer le son"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 16 : Code Vidéo
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Code : Vidéo de Fond"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """// JavaScript (Fatoumata)
function onYouTubeIframeAPIReady() {
    player = new YT.Player('background-video', {
        videoId: 'ZsJz2TJAPjw',
        playerVars: {
            'autoplay': 1,
            'mute': 1,
            'loop': 1,
        },
        events: {
            'onReady': function(event) {
                event.target.playVideo();
            }
        }
    });
}

function toggleAudio() {
    if (isMuted) {
        player.unMute();
        player.setVolume(70);
    } else {
        player.mute();
    }
}"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 17-18 : Optimisation Performance
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Optimisation Performance"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"L\'application est lente. Comment optimiser ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Système de cache JSON"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Indicateurs de chargement visuels"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Lazy loading des données"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 18 : Code Cache
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Code : Système de Cache"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """# cache_manager.py (Hector)
def get_cached_films() -> List[Dict]:
    \"\"\"Charge depuis le cache.\"\"\"
    if CACHE_FILE.exists():
        with open(CACHE_FILE, "r") as f:
            films = json.load(f)
        return films
    return []

def cache_films(films: List[Dict]) -> None:
    \"\"\"Sauvegarde dans le cache.\"\"\"
    with open(CACHE_FILE, "w") as f:
        json.dump(films, f, indent=2)"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(10)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 19-20 : Détection d'Émotions
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Détection d'Émotions Faciales"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Comment ajouter la détection d\'émotions via webcam ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ DeepFace (modèle pré-entraîné)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ OpenCV pour détection visage"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Mapping émotions DeepFace → nos émotions"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 20 : Code DeepFace
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Code : Détection d'Émotions"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💻 CODE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    code_text = """# emotion_detection.py (Hector)
from deepface import DeepFace
import cv2

def detecter_emotion_image(image_data: bytes):
    \"\"\"Détecte l'émotion avec DeepFace.\"\"\"
    img = cv2.imdecode(image_data, cv2.IMREAD_COLOR)
    
    # Détecter visage
    faces = face_cascade.detectMultiScale(gray)
    
    if len(faces) > 0:
        # Analyser avec DeepFace
        result = DeepFace.analyze(
            img_path=tmp_path,
            actions=['emotion']
        )
        emotion = EMOTION_MAPPING.get(
            result['dominant_emotion'], 
            "neutre"
        )
        return {"emotion": emotion}"""
    
    p = tf.add_paragraph()
    p.text = code_text
    p.font.size = Pt(9)
    p.font.name = "Courier New"
    p.level = 1
    
    # ============================================
    # SLIDE 21-22 : Évaluation
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Évaluation & Notebooks"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🤔 PROMPT :"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = '"Comment évaluer les performances de tous nos modèles ?"'
    p.font.size = Pt(14)
    p.font.italic = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 RÉPONSE :"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Notebook Jupyter complet"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Analyse qualité dataset"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Évaluation recommandations"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Matrices de confusion"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Visualisations"
    p.font.size = Pt(14)
    p.level = 1
    
    # Slide 22 : Résultats Évaluation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Résultats de l'Évaluation"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "📊 MÉTRIQUES CLÉS :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Dataset : 4,803 films (99.4% avec genres)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Recommandations : 8/8 émotions (100%)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Note moyenne recommandations : 8.50/10"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Amélioration : +2.41 points vs moyenne globale"
    p.font.size = Pt(14)
    p.level = 1
    
    # ============================================
    # SLIDE 23 : Résultats & Métriques
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Résultats & Performances"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "📊 PERFORMANCES :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\n✅ Dataset :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 4,803 films chargés"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• 99.4% avec genres"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• 98.7% avec notes valides"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\n✅ Système de Recommandation :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 8/8 émotions couvertes (100%)"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Note moyenne : 8.50/10"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• +2.41 points vs moyenne globale"
    p.font.size = Pt(14)
    p.level = 2
    
    # ============================================
    # SLIDE 24 : Démo Live
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "🎬 Démonstration Live"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "1. Recherche par titre"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "   Tapez 'The Matrix' → Voir les détails"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n2. Recommandation par émotion"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "   Choisissez 'stressé' → 20 films adaptés"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n3. Détection faciale (bonus)"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "   Activez webcam → IA détecte émotion"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n4. Vidéo de fond"
    p.font.size = Pt(16)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "   Ambiance cinématographique + Son"
    p.font.size = Pt(14)
    p.level = 1
    
    # ============================================
    # SLIDE 25 : Technologies
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Technologies Utilisées"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🔧 STACK TECHNIQUE :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\nBack-end : Python 3.12, Flask, Pandas"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nIA & ML : TextBlob, DeepFace, TensorFlow, OpenCV"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nFront-end : HTML5, CSS3, JavaScript, YouTube API"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nDonnées : TMDB (4,803 films), Cache JSON"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\nÉvaluation : Jupyter, Matplotlib, Seaborn"
    p.font.size = Pt(14)
    p.level = 1
    
    # ============================================
    # SLIDE 26 : Défis Relevés
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Défis Techniques Relevés"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "💪 DÉFIS RÉSOLUS :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "1. ✅ Autoplay vidéo → YouTube Iframe API"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. ✅ Conflits dépendances → Gestion versions"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. ✅ Performance → Système de cache"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. ✅ Évaluation modèles → Notebook complet"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "5. ✅ Erreurs Jupyter → Réorganisation imports"
    p.font.size = Pt(14)
    p.level = 1
    
    # ============================================
    # SLIDE 27 : Améliorations Futures
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Améliorations Futures"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🚀 ROADMAP :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\nCourt terme :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Fine-tuning TextBlob"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Optimisation DeepFace"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Tests utilisateurs"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\nMoyen terme :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Filtrage collaboratif"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Personnalisation utilisateur"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\nLong terme :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Déploiement cloud"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Application mobile"
    p.font.size = Pt(14)
    p.level = 2
    
    # ============================================
    # SLIDE 28 : Conclusion
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Conclusion"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🎯 CE QUE NOUS AVONS RÉALISÉ :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "✅ Plateforme web complète et fonctionnelle"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Système de recommandation intelligent"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Analyse de sentiments intégrée"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Détection d'émotions faciales (bonus)"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Interface utilisateur moderne"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "✅ Évaluation complète des modèles"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "\n💡 VALEUR AJOUTÉE :"
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "🎬 Pour les utilisateurs : Gain de temps, découvertes personnalisées"
    p.font.size = Pt(14)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "📊 Pour l'équipe : Maîtrise IA, expérience web, collaboration"
    p.font.size = Pt(14)
    p.level = 1
    
    # ============================================
    # SLIDE 29 : Remerciements
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Remerciements"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "🙏 REMERCIEMENTS :"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\nMerci à notre encadrante pour :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Son accompagnement"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Ses conseils précieux"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Sa flexibilité"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "\n📚 RESSOURCES :"
    p.font.size = Pt(16)
    p.font.bold = True
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Dataset : TMDB 5000 Movies"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Modèles : TextBlob, DeepFace"
    p.font.size = Pt(14)
    p.level = 2
    
    p = tf.add_paragraph()
    p.text = "• Documentation : Flask, TensorFlow"
    p.font.size = Pt(14)
    p.level = 2
    
    # ============================================
    # SLIDE 30 : Questions
    # ============================================
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Questions ?"
    title.text_frame.paragraphs[0].font.color.rgb = ACCENT_COLOR
    
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "❓ QUESTIONS & RÉPONSES"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\n\nContact :"
    p.font.size = Pt(18)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Gémima ONDELE POURU"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Fatoumata BAH"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "Hector KOMBOU"
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.CENTER
    
    p = tf.add_paragraph()
    p.text = "\nCode source : Disponible sur GitHub"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.CENTER
    
    # Sauvegarder la présentation
    output_path = Path(__file__).parent.parent / "MovieMood_Presentation.pptx"
    prs.save(str(output_path))
    
    print(f"✅ Présentation créée avec succès !")
    print(f"📁 Fichier : {output_path}")
    print(f"📊 Nombre de slides : {len(prs.slides)}")
    
    return output_path

if __name__ == "__main__":
    try:
        output_path = create_presentation()
        print(f"\n🎉 La présentation est prête dans : {output_path}")
        print("💡 Vous pouvez maintenant l'ouvrir dans PowerPoint ou l'importer dans Canva !")
    except ImportError:
        print("❌ Erreur : python-pptx n'est pas installé.")
        print("💡 Installez-le avec : pip install python-pptx")
    except Exception as e:
        print(f"❌ Erreur lors de la création : {e}")
        import traceback
        traceback.print_exc()

